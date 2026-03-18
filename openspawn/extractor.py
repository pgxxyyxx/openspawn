from __future__ import annotations

import csv
import json
import subprocess
from pathlib import Path

from .types import ExtractionResult, FileEntry, Limits


def extract_file(path: Path, entry: FileEntry, limits: Limits) -> ExtractionResult:
    if entry.ext in {".txt", ".md", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".toml", ".ini"}:
        return _extract_text(path, entry, limits)
    if entry.ext in {".csv", ".tsv"}:
        return _extract_csv(path, entry)
    if entry.ext == ".ipynb":
        return _extract_notebook(path, limits)
    if entry.ext in {".xlsx", ".xls"}:
        return _extract_excel(path)
    if entry.ext == ".pdf":
        return _extract_pdf(path, limits)
    if entry.ext == ".docx":
        return _extract_docx(path, limits)
    if entry.ext == ".pptx":
        return _extract_pptx(path, limits)
    return ExtractionResult(
        status="skipped_unsupported",
        parser="",
        bytes_read=0,
        content_text="",
        structured_summary={"kind": "unsupported"},
        warnings=["Unsupported file type in Stage 0."],
    )


def _extract_text(path: Path, entry: FileEntry, limits: Limits) -> ExtractionResult:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except PermissionError:
        return _error("error_permission", "text", "Permission denied.")
    except OSError as exc:
        return _error("error_parse", "text", str(exc))
    if len(text.encode("utf-8")) <= limits.max_text_bytes:
        sample = text
        status = "read_full"
        warnings: list[str] = []
    else:
        sample = _sample_text(text, limits.max_text_bytes)
        status = "read_partial"
        warnings = ["Text was sampled because it exceeded the size limit."]
    return ExtractionResult(
        status=status,
        parser="text",
        bytes_read=len(sample.encode("utf-8")),
        content_text=sample,
        structured_summary={"kind": "text", "line_count": text.count("\n") + 1},
        warnings=warnings,
    )


def _extract_csv(path: Path, entry: FileEntry) -> ExtractionResult:
    delimiter = "\t" if entry.ext == ".tsv" else ","
    try:
        with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            reader = csv.reader(handle, delimiter=delimiter)
            headers = next(reader, [])
            sample_rows: list[list[str]] = []
            row_count = 0
            for row in reader:
                row_count += 1
                if len(sample_rows) < 3:
                    sample_rows.append(row[:10])
                if row_count >= 5000:
                    break
    except PermissionError:
        return _error("error_permission", "csv", "Permission denied.")
    except OSError as exc:
        return _error("error_parse", "csv", str(exc))
    content_lines = [f"Columns: {', '.join(headers[:20])}", f"Rows counted: {row_count}"]
    for row in sample_rows:
        content_lines.append(f"Sample: {row}")
    return ExtractionResult(
        status="read_partial",
        parser="csv",
        bytes_read=sum(len(line.encode("utf-8")) for line in content_lines),
        content_text="\n".join(content_lines),
        structured_summary={"kind": "table", "columns": headers, "row_count": row_count},
    )


def _extract_notebook(path: Path, limits: Limits) -> ExtractionResult:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except PermissionError:
        return _error("error_permission", "notebook", "Permission denied.")
    except (OSError, json.JSONDecodeError) as exc:
        return _error("error_parse", "notebook", str(exc))
    cells = data.get("cells", [])
    parts: list[str] = []
    for cell in cells:
        label = "CODE" if cell.get("cell_type") == "code" else "MD"
        chunk = "".join(cell.get("source", []))
        if chunk.strip():
            parts.append(f"[{label}]\n{chunk[:500]}")
        if sum(len(part) for part in parts) >= limits.max_text_bytes:
            break
    text = "\n\n".join(parts)[: limits.max_text_bytes]
    return ExtractionResult(
        status="read_partial" if len(text) >= limits.max_text_bytes else "read_full",
        parser="ipynb",
        bytes_read=len(text.encode("utf-8")),
        content_text=text,
        structured_summary={"kind": "notebook", "cell_count": len(cells)},
    )


def _extract_excel(path: Path) -> ExtractionResult:
    try:
        import openpyxl
    except ImportError:
        return _error("error_parse", "openpyxl", "openpyxl is not installed.")
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except PermissionError:
        return _error("error_permission", "openpyxl", "Permission denied.")
    except Exception as exc:
        return _error("error_parse", "openpyxl", str(exc))
    lines: list[str] = []
    sheets: list[dict[str, object]] = []
    try:
        for sheet_name in workbook.sheetnames[:10]:
            sheet = workbook[sheet_name]
            rows = list(sheet.iter_rows(max_row=5, values_only=True))
            headers = [str(cell) if cell is not None else "" for cell in rows[0]] if rows else []
            sheets.append({"name": sheet_name, "rows": sheet.max_row, "columns": sheet.max_column})
            lines.append(f"Sheet: {sheet_name} ({sheet.max_row} rows x {sheet.max_column} columns)")
            if headers:
                lines.append(f"Headers: {', '.join(headers[:15])}")
    finally:
        workbook.close()
    text = "\n".join(lines)
    return ExtractionResult(
        status="read_partial",
        parser="openpyxl",
        bytes_read=len(text.encode("utf-8")),
        content_text=text,
        structured_summary={"kind": "workbook", "sheets": sheets},
    )


def _extract_pdf(path: Path, limits: Limits) -> ExtractionResult:
    try:
        from pypdf import PdfReader
    except ImportError:
        PdfReader = None

    if PdfReader is not None:
        try:
            reader = PdfReader(str(path))
            parts: list[str] = []
            pages_read = 0
            for page in reader.pages[: limits.max_pdf_pages]:
                pages_read += 1
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text)
                if sum(len(part.encode("utf-8")) for part in parts) >= limits.max_text_bytes:
                    break
            combined = "\n\n".join(parts)[: limits.max_text_bytes]
            warnings: list[str] = []
            status = "read_full"
            if len(reader.pages) > pages_read or len(combined.encode("utf-8")) >= limits.max_text_bytes:
                status = "read_partial"
                warnings.append("PDF text was capped at the extraction limit.")
            if combined.strip():
                return ExtractionResult(
                    status=status,
                    parser="pypdf",
                    bytes_read=len(combined.encode("utf-8")),
                    content_text=combined,
                    structured_summary={"kind": "pdf", "page_count": len(reader.pages)},
                    warnings=warnings,
                )
        except PermissionError:
            return _error("error_permission", "pypdf", "Permission denied.")
        except Exception as exc:
            # Fall through to CLI tool fallback before giving up completely.
            fallback_warning = str(exc)
        else:
            fallback_warning = ""
    else:
        fallback_warning = "pypdf is not installed."

    try:
        result = subprocess.run(
            ["pdftotext", "-l", str(limits.max_pdf_pages), str(path), "-"],
            capture_output=True,
            text=True,
            timeout=10,
            check=False,
        )
    except FileNotFoundError:
        return _error("error_parse", "pdftotext", "pdftotext is not installed.")
    except subprocess.TimeoutExpired:
        return _error("error_timeout", "pdftotext", "PDF extraction timed out.")
    if result.returncode != 0 or not result.stdout.strip():
        return ExtractionResult(
            status="read_metadata_only",
            parser="pdftotext" if "pdftotext" in result.args[0] else "pdf",
            bytes_read=0,
            content_text="",
            structured_summary={"kind": "pdf", "page_count": None},
            warnings=["PDF text extraction failed; metadata only available."],
            error=result.stderr.strip() or fallback_warning or "PDF extraction failed.",
        )
    text = result.stdout[: limits.max_text_bytes]
    warnings = []
    status = "read_full"
    if len(result.stdout) > limits.max_text_bytes:
        status = "read_partial"
        warnings.append("PDF text was capped at the extraction limit.")
    return ExtractionResult(
        status=status,
        parser="pdftotext",
        bytes_read=len(text.encode("utf-8")),
        content_text=text,
        structured_summary={"kind": "pdf", "page_count": limits.max_pdf_pages},
        warnings=warnings,
    )


def _extract_docx(path: Path, limits: Limits) -> ExtractionResult:
    try:
        from docx import Document
    except ImportError:
        return _error("error_parse", "python-docx", "python-docx is not installed.")
    try:
        doc = Document(str(path))
    except PermissionError:
        return _error("error_permission", "python-docx", "Permission denied.")
    except Exception as exc:
        return _error("error_parse", "python-docx", str(exc))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    combined = "\n".join(parts)[: limits.max_text_bytes]
    status = "read_partial" if len("\n".join(parts).encode("utf-8")) > limits.max_text_bytes else "read_full"
    return ExtractionResult(
        status=status,
        parser="python-docx",
        bytes_read=len(combined.encode("utf-8")),
        content_text=combined,
        structured_summary={"kind": "document", "paragraph_count": len(doc.paragraphs)},
    )


def _extract_pptx(path: Path, limits: Limits) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return _error("error_parse", "python-pptx", "python-pptx is not installed.")
    try:
        prs = Presentation(str(path))
    except PermissionError:
        return _error("error_permission", "python-pptx", "Permission denied.")
    except Exception as exc:
        return _error("error_parse", "python-pptx", str(exc))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    combined = "\n\n".join(parts)[: limits.max_text_bytes]
    status = "read_partial" if len("\n\n".join(parts).encode("utf-8")) > limits.max_text_bytes else "read_full"
    return ExtractionResult(
        status=status,
        parser="python-pptx",
        bytes_read=len(combined.encode("utf-8")),
        content_text=combined,
        structured_summary={"kind": "presentation", "slide_count": len(prs.slides)},
    )


def _sample_text(text: str, max_bytes: int) -> str:
    target_chars = max(max_bytes // 3, 200)
    if len(text) <= target_chars * 3:
        return text[:max_bytes]
    head = text[:target_chars]
    middle_start = max(0, len(text) // 2 - target_chars // 2)
    middle = text[middle_start : middle_start + target_chars]
    tail = text[-target_chars:]
    return "\n...\n".join([head, middle, tail])[:max_bytes]


def _error(status: str, parser: str, message: str) -> ExtractionResult:
    return ExtractionResult(
        status=status,
        parser=parser,
        bytes_read=0,
        content_text="",
        structured_summary={"kind": "error"},
        warnings=[],
        error=message,
    )
